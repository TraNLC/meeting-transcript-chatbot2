from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Define colors
PRIMARY_COLOR = RGBColor(102, 126, 234)  # #667eea
SECONDARY_COLOR = RGBColor(118, 75, 162)  # #764ba2
TEXT_COLOR = RGBColor(31, 41, 55)  # #1f2937
GRAY_COLOR = RGBColor(107, 114, 128)  # #6b7280

def add_title_slide(prs):
    """Slide 1: Title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Meeting Transcript Chatbot"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(14), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Hệ thống phân tích cuộc họp thông minh với AI"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = GRAY_COLOR
    subtitle_para.alignment = PP_ALIGN.CENTER

def add_feature_slide(prs, title, icon, goal, stats, workflow_steps, tech_items, highlight):
    """Add a feature slide with workflow"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Header
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    header_frame = header_box.text_frame
    header_frame.text = f"{icon} {title}"
    header_para = header_frame.paragraphs[0]
    header_para.font.size = Pt(40)
    header_para.font.bold = True
    header_para.font.color.rgb = PRIMARY_COLOR
    
    # Goal box
    goal_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(15), Inches(0.6))
    goal_frame = goal_box.text_frame
    goal_frame.text = f"🎯 Mục tiêu: {goal}"
    goal_para = goal_frame.paragraphs[0]
    goal_para.font.size = Pt(18)
    goal_para.font.color.rgb = TEXT_COLOR
    
    # Stats (3 columns)
    stat_y = 2.0
    stat_width = 4.5
    for i, (stat_value, stat_label) in enumerate(stats):
        stat_x = 0.5 + i * 5
        stat_box = slide.shapes.add_textbox(Inches(stat_x), Inches(stat_y), Inches(stat_width), Inches(0.8))
        stat_frame = stat_box.text_frame
        
        # Value
        stat_frame.text = str(stat_value)
        value_para = stat_frame.paragraphs[0]
        value_para.font.size = Pt(36)
        value_para.font.bold = True
        value_para.font.color.rgb = PRIMARY_COLOR
        value_para.alignment = PP_ALIGN.CENTER
        
        # Label
        label_para = stat_frame.add_paragraph()
        label_para.text = stat_label
        label_para.font.size = Pt(14)
        label_para.font.color.rgb = GRAY_COLOR
        label_para.alignment = PP_ALIGN.CENTER
    
    # Workflow title
    workflow_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(7), Inches(0.4))
    workflow_title_frame = workflow_title_box.text_frame
    workflow_title_frame.text = "📋 Workflow"
    workflow_title_para = workflow_title_frame.paragraphs[0]
    workflow_title_para.font.size = Pt(22)
    workflow_title_para.font.bold = True
    workflow_title_para.font.color.rgb = TEXT_COLOR
    
    # Workflow steps (left column)
    workflow_y = 3.8
    for i, (step_title, step_desc) in enumerate(workflow_steps):
        step_box = slide.shapes.add_textbox(Inches(0.5), Inches(workflow_y), Inches(7), Inches(0.6))
        step_frame = step_box.text_frame
        
        # Step number and title
        step_frame.text = f"{i+1}. {step_title}"
        step_para = step_frame.paragraphs[0]
        step_para.font.size = Pt(14)
        step_para.font.bold = True
        step_para.font.color.rgb = TEXT_COLOR
        
        # Description
        desc_para = step_frame.add_paragraph()
        desc_para.text = step_desc
        desc_para.font.size = Pt(11)
        desc_para.font.color.rgb = GRAY_COLOR
        desc_para.level = 1
        
        workflow_y += 0.65
    
    # Tech stack title
    tech_title_box = slide.shapes.add_textbox(Inches(8.5), Inches(3.2), Inches(7), Inches(0.4))
    tech_title_frame = tech_title_box.text_frame
    tech_title_frame.text = "🔧 Công Nghệ"
    tech_title_para = tech_title_frame.paragraphs[0]
    tech_title_para.font.size = Pt(22)
    tech_title_para.font.bold = True
    tech_title_para.font.color.rgb = TEXT_COLOR
    
    # Tech items (right column)
    tech_y = 3.8
    for tech_name, tech_purpose in tech_items:
        tech_box = slide.shapes.add_textbox(Inches(8.5), Inches(tech_y), Inches(7), Inches(0.6))
        tech_frame = tech_box.text_frame
        
        # Tech name
        tech_frame.text = f"• {tech_name}"
        tech_para = tech_frame.paragraphs[0]
        tech_para.font.size = Pt(13)
        tech_para.font.bold = True
        tech_para.font.color.rgb = PRIMARY_COLOR
        
        # Purpose
        purpose_para = tech_frame.add_paragraph()
        purpose_para.text = tech_purpose
        purpose_para.font.size = Pt(10)
        purpose_para.font.color.rgb = GRAY_COLOR
        purpose_para.level = 1
        
        tech_y += 0.65
    
    # Highlight box at bottom
    highlight_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.8), Inches(15), Inches(0.8))
    highlight_frame = highlight_box.text_frame
    highlight_frame.text = f"💡 {highlight}"
    highlight_para = highlight_frame.paragraphs[0]
    highlight_para.font.size = Pt(14)
    highlight_para.font.italic = True
    highlight_para.font.color.rgb = PRIMARY_COLOR

# Slide 1: Title
add_title_slide(prs)

# Slide 2: Recording & Transcription
add_feature_slide(
    prs,
    title="Ghi Âm & Chuyển Văn Bản",
    icon="🎙️",
    goal="Ghi âm cuộc họp trực tiếp và tự động chuyển thành văn bản với speaker diarization",
    stats=[
        ("92%", "Độ chính xác"),
        ("50+", "Ngôn ngữ"),
        ("3x", "Realtime")
    ],
    workflow_steps=[
        ("Khởi động ghi âm", "WebRTC API capture audio từ browser/microphone"),
        ("Streaming & Buffer", "Audio chunks được buffer và gửi real-time qua WebSocket"),
        ("Transcription (WhisperX)", "WhisperX Medium model chuyển audio → text với timestamp"),
        ("Speaker Diarization", "PyAnnote 3.1 phân biệt người nói (Speaker 1, 2, 3...)"),
        ("Lưu & Hiển thị", "Lưu transcript vào DB, hiển thị real-time trên UI")
    ],
    tech_items=[
        ("WhisperX (OpenAI Whisper)", "Speech-to-Text với độ chính xác 92%, hỗ trợ 50+ ngôn ngữ"),
        ("PyAnnote Audio 3.1", "Speaker Diarization - phân biệt người nói"),
        ("WebRTC + WebSocket", "Real-time audio streaming từ browser"),
        ("FFmpeg", "Audio processing và format conversion")
    ],
    highlight="Sử dụng DiarizationPipeline wrapper để đảm bảo output format tương thích"
)

# Slide 3: Upload & Analysis
add_feature_slide(
    prs,
    title="Upload & Phân Tích",
    icon="📤",
    goal="Upload file và nhận phân tích AI toàn diện về cuộc họp",
    stats=[
        ("100MB", "Max Size"),
        ("5", "File Formats"),
        ("30s", "Avg Time")
    ],
    workflow_steps=[
        ("Upload File", "Hỗ trợ TXT, DOCX, MP3, WAV, MP4 (max 100MB)"),
        ("File Validation", "Kiểm tra định dạng, kích thước, virus scan"),
        ("Content Extraction", "Audio → WhisperX, Text → Direct, Video → FFmpeg extract"),
        ("AI Analysis (Gemini 2.5)", "Tóm tắt, trích xuất chủ đề, action items, quyết định"),
        ("Vectorization & Storage", "Tạo embeddings và lưu vào ChromaDB + PostgreSQL")
    ],
    tech_items=[
        ("Google Gemini 2.5 Flash", "LLM chính cho phân tích, tóm tắt, trích xuất thông tin"),
        ("Prompt Engineering", "6 kỹ thuật: Few-shot, Chain-of-Thought, Role-based"),
        ("python-docx", "Đọc/ghi file Word với formatting"),
        ("Input Sanitization", "Validate, clean, truncate input để tránh injection")
    ],
    highlight="Sử dụng Structured Output với JSON schema để đảm bảo format nhất quán"
)

# Slide 4: RAG Chat
add_feature_slide(
    prs,
    title="Chat với AI (RAG)",
    icon="💬",
    goal="Hỏi đáp thông minh về cuộc họp với Retrieval-Augmented Generation",
    stats=[
        ("384", "Embedding Dims"),
        ("5", "Top-K Results"),
        ("10", "Chat History")
    ],
    workflow_steps=[
        ("User Query", "Người dùng đặt câu hỏi về cuộc họp"),
        ("Query Embedding", "Chuyển câu hỏi thành vector (sentence-transformers)"),
        ("Similarity Search", "Tìm top-k chunks liên quan nhất (cosine similarity)"),
        ("Context Assembly", "Kết hợp retrieved chunks + conversation history"),
        ("LLM Generation", "Gemini tạo câu trả lời dựa trên context"),
        ("Memory Update", "Lưu Q&A vào conversation memory")
    ],
    tech_items=[
        ("ChromaDB", "Vector Database - lưu trữ và tìm kiếm embeddings"),
        ("LangChain", "Framework RAG: RetrievalQA, ConversationBufferMemory"),
        ("sentence-transformers", "Tạo embeddings (all-MiniLM-L6-v2, 384 dimensions)"),
        ("Conversation Memory", "Lưu lịch sử chat để xử lý follow-up questions")
    ],
    highlight="Kết hợp RAG + Conversation Memory để chatbot hiểu context và câu hỏi tiếp theo"
)

# Slide 5: Semantic Search
add_feature_slide(
    prs,
    title="Tìm Kiếm Ngữ Nghĩa",
    icon="🔍",
    goal="Tìm kiếm cuộc họp theo ý nghĩa, không chỉ từ khóa",
    stats=[
        ("384", "Vector Dims"),
        ("95%", "Accuracy"),
        ("100ms", "Search Time")
    ],
    workflow_steps=[
        ("Search Query", "User nhập query (VD: 'cuộc họp về marketing')"),
        ("Query Embedding", "Chuyển query thành vector 384 chiều"),
        ("Vector Search + Filters", "ChromaDB tìm kiếm với metadata filters"),
        ("Ranking & Scoring", "Sắp xếp theo similarity score (0-1)"),
        ("Display Results", "Hiển thị kết quả với snippet và highlight")
    ],
    tech_items=[
        ("ChromaDB Advanced", "Metadata filtering, similarity search, collection management"),
        ("Cosine Similarity", "Đo độ tương đồng giữa vectors (range: 0-1)"),
        ("Metadata Filters", "Lọc theo ngày, ngôn ngữ, loại cuộc họp, người tham gia"),
        ("Analytics Dashboard", "Thống kê database: số lượng, phân bố, trends")
    ],
    highlight="Kết hợp vector search với metadata filtering để tìm kiếm chính xác và nhanh"
)

# Slide 6: History & Export
add_feature_slide(
    prs,
    title="Lịch Sử & Xuất File",
    icon="📊",
    goal="Quản lý lịch sử cuộc họp và xuất kết quả ra nhiều định dạng",
    stats=[
        ("10", "Items/Page"),
        ("2", "Export Formats"),
        ("500+", "Meetings Stored")
    ],
    workflow_steps=[
        ("Load History", "Fetch từ PostgreSQL với pagination (10 items/page)"),
        ("Display Cards", "Hiển thị meeting cards với preview, date, participants"),
        ("Select Meeting", "Click vào card để xem chi tiết đầy đủ"),
        ("Export Options", "Chọn format: TXT (plain), DOCX (formatted)"),
        ("Generate & Download", "Tạo file với formatting và trigger download")
    ],
    tech_items=[
        ("PostgreSQL", "Lưu trữ metadata, transcript, analysis results"),
        ("python-docx", "Tạo file Word với headers, bullets, tables, formatting"),
        ("Pagination", "Load dữ liệu theo batch để tối ưu performance"),
        ("Soft Delete", "Đánh dấu xóa thay vì xóa vĩnh viễn (data recovery)")
    ],
    highlight="Dual storage (PostgreSQL + ChromaDB) để tối ưu cả structured data và vector search"
)

# Save presentation
prs.save('Meeting_Transcript_Chatbot_Presentation.pptx')
print("Da tao file: Meeting_Transcript_Chatbot_Presentation.pptx")
