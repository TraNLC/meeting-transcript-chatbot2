from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Define colors
PRIMARY_COLOR = RGBColor(102, 126, 234)  # #667eea
SECONDARY_COLOR = RGBColor(118, 75, 162)  # #764ba2
TEXT_COLOR = RGBColor(31, 41, 55)  # #1f2937
GRAY_COLOR = RGBColor(107, 114, 128)  # #6b7280
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

def add_table_of_contents_slide(prs):
    """Create Table of Contents slide similar to the image"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(20, 20, 30)
    
    # Title "TABLE OF CONTENTS" with border
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(8), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = "TABLE OF CONTENTS"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(56)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.LEFT
    
    # Add border around title (using shape)
    border_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.7), Inches(0.7),
        Inches(8.2), Inches(1.4)
    )
    border_shape.fill.background()
    border_shape.line.color.rgb = PRIMARY_COLOR
    border_shape.line.width = Pt(3)
    
    # Content items
    content_items = [
        ("01", "Giới Thiệu Dự Án"),
        ("02", "Kiến Trúc Hệ Thống"),
        ("03", "Tính Năng Chính"),
        ("04", "🎙️ Ghi Âm & Transcription"),
        ("05", "📤 Upload & Phân Tích AI"),
        ("06", "💬 Chat với AI (RAG)"),
        ("07", "🔍 Tìm Kiếm Ngữ Nghĩa"),
        ("08", "Demo Trực Tiếp"),
        ("09", "Q&A")
    ]
    
    start_y = 2.5
    line_height = 0.55
    
    for i, (number, text) in enumerate(content_items):
        y_pos = start_y + (i * line_height)
        
        # Number
        num_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(1), Inches(0.5))
        num_frame = num_box.text_frame
        num_frame.text = number
        num_para = num_frame.paragraphs[0]
        num_para.font.size = Pt(32)
        num_para.font.bold = True
        num_para.font.color.rgb = WHITE
        
        # Bullet point
        bullet_box = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(0.3), Inches(0.5))
        bullet_frame = bullet_box.text_frame
        bullet_frame.text = "•"
        bullet_para = bullet_frame.paragraphs[0]
        bullet_para.font.size = Pt(28)
        bullet_para.font.color.rgb = PRIMARY_COLOR
        
        # Text
        text_box = slide.shapes.add_textbox(Inches(2.5), Inches(y_pos), Inches(6), Inches(0.5))
        text_frame = text_box.text_frame
        text_frame.text = text
        text_para = text_frame.paragraphs[0]
        text_para.font.size = Pt(28)
        text_para.font.color.rgb = WHITE
        
        # Highlight main features (04-07)
        if i >= 3 and i <= 6:
            text_para.font.bold = True
            text_para.font.color.rgb = RGBColor(100, 200, 255)  # Light blue

def add_title_slide(prs):
    """Slide 1: Title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Gradient background (dark)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(20, 20, 40)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Meeting Transcript Chatbot"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(14), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Hệ thống phân tích cuộc họp thông minh với AI"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(150, 150, 200)
    subtitle_para.alignment = PP_ALIGN.CENTER

def add_feature_slide(prs, title, icon, goal, stats, workflow_steps, tech_items, highlight):
    """Add a feature slide with workflow"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Dark background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(20, 20, 40)
    
    # Header with icon
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    header_frame = header_box.text_frame
    header_frame.text = f"{icon} {title}"
    header_para = header_frame.paragraphs[0]
    header_para.font.size = Pt(40)
    header_para.font.bold = True
    header_para.font.color.rgb = WHITE
    
    # Goal box
    goal_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(15), Inches(0.6))
    goal_frame = goal_box.text_frame
    goal_frame.text = f"🎯 {goal}"
    goal_para = goal_frame.paragraphs[0]
    goal_para.font.size = Pt(18)
    goal_para.font.color.rgb = RGBColor(200, 200, 255)
    
    # Stats (3 columns)
    stat_y = 2.0
    stat_width = 4.5
    for i, (stat_value, stat_label) in enumerate(stats):
        stat_x = 0.5 + i * 5
        
        # Stat card background
        stat_bg = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(stat_x), Inches(stat_y),
            Inches(stat_width), Inches(0.8)
        )
        stat_bg.fill.solid()
        stat_bg.fill.fore_color.rgb = RGBColor(40, 40, 60)
        stat_bg.line.color.rgb = PRIMARY_COLOR
        stat_bg.line.width = Pt(2)
        
        stat_box = slide.shapes.add_textbox(Inches(stat_x), Inches(stat_y), Inches(stat_width), Inches(0.8))
        stat_frame = stat_box.text_frame
        
        # Value
        stat_frame.text = str(stat_value)
        value_para = stat_frame.paragraphs[0]
        value_para.font.size = Pt(36)
        value_para.font.bold = True
        value_para.font.color.rgb = PRIMARY_COLOR
        value_para.alignment = PP_ALIGN.CENTER
        
        # Label
        label_para = stat_frame.add_paragraph()
        label_para.text = stat_label
        label_para.font.size = Pt(14)
        label_para.font.color.rgb = RGBColor(180, 180, 200)
        label_para.alignment = PP_ALIGN.CENTER
    
    # Workflow title
    workflow_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(7), Inches(0.4))
    workflow_title_frame = workflow_title_box.text_frame
    workflow_title_frame.text = "📋 WORKFLOW"
    workflow_title_para = workflow_title_frame.paragraphs[0]
    workflow_title_para.font.size = Pt(22)
    workflow_title_para.font.bold = True
    workflow_title_para.font.color.rgb = WHITE
    
    # Workflow steps (left column)
    workflow_y = 3.8
    for i, (step_title, step_desc) in enumerate(workflow_steps):
        # Step background
        step_bg = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.5), Inches(workflow_y - 0.05),
            Inches(7), Inches(0.6)
        )
        step_bg.fill.solid()
        step_bg.fill.fore_color.rgb = RGBColor(30, 30, 50)
        step_bg.line.color.rgb = PRIMARY_COLOR
        step_bg.line.width = Pt(1)
        
        step_box = slide.shapes.add_textbox(Inches(0.7), Inches(workflow_y), Inches(6.5), Inches(0.6))
        step_frame = step_box.text_frame
        
        # Step number and title
        step_frame.text = f"{i+1}. {step_title}"
        step_para = step_frame.paragraphs[0]
        step_para.font.size = Pt(14)
        step_para.font.bold = True
        step_para.font.color.rgb = WHITE
        
        # Description
        desc_para = step_frame.add_paragraph()
        desc_para.text = step_desc
        desc_para.font.size = Pt(11)
        desc_para.font.color.rgb = RGBColor(160, 160, 180)
        desc_para.level = 1
        
        workflow_y += 0.65
    
    # Tech stack title
    tech_title_box = slide.shapes.add_textbox(Inches(8.5), Inches(3.2), Inches(7), Inches(0.4))
    tech_title_frame = tech_title_box.text_frame
    tech_title_frame.text = "🔧 CÔNG NGHỆ"
    tech_title_para = tech_title_frame.paragraphs[0]
    tech_title_para.font.size = Pt(22)
    tech_title_para.font.bold = True
    tech_title_para.font.color.rgb = WHITE
    
    # Tech items (right column)
    tech_y = 3.8
    for tech_name, tech_purpose in tech_items:
        tech_box = slide.shapes.add_textbox(Inches(8.5), Inches(tech_y), Inches(7), Inches(0.6))
        tech_frame = tech_box.text_frame
        
        # Tech name
        tech_frame.text = f"• {tech_name}"
        tech_para = tech_frame.paragraphs[0]
        tech_para.font.size = Pt(13)
        tech_para.font.bold = True
        tech_para.font.color.rgb = PRIMARY_COLOR
        
        # Purpose
        purpose_para = tech_frame.add_paragraph()
        purpose_para.text = tech_purpose
        purpose_para.font.size = Pt(10)
        purpose_para.font.color.rgb = RGBColor(160, 160, 180)
        purpose_para.level = 1
        
        tech_y += 0.65
    
    # Highlight box at bottom
    highlight_bg = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(7.8),
        Inches(15), Inches(0.8)
    )
    highlight_bg.fill.solid()
    highlight_bg.fill.fore_color.rgb = RGBColor(40, 40, 60)
    highlight_bg.line.color.rgb = RGBColor(255, 200, 0)
    highlight_bg.line.width = Pt(2)
    
    highlight_box = slide.shapes.add_textbox(Inches(0.7), Inches(7.9), Inches(14.5), Inches(0.6))
    highlight_frame = highlight_box.text_frame
    highlight_frame.text = f"💡 {highlight}"
    highlight_para = highlight_frame.paragraphs[0]
    highlight_para.font.size = Pt(14)
    highlight_para.font.italic = True
    highlight_para.font.color.rgb = RGBColor(255, 220, 100)

# Slide 1: Title
add_title_slide(prs)

# Slide 2: Table of Contents
add_table_of_contents_slide(prs)

# Slide 3: Recording & Transcription
add_feature_slide(
    prs,
    title="Ghi Âm & Chuyển Văn Bản",
    icon="🎙️",
    goal="Ghi âm cuộc họp và tự động chuyển thành văn bản với speaker diarization",
    stats=[
        ("92%", "Độ chính xác"),
        ("50+", "Ngôn ngữ"),
        ("3x", "Realtime")
    ],
    workflow_steps=[
        ("Khởi động ghi âm", "WebRTC API capture audio từ browser/mic"),
        ("Streaming & Buffer", "Audio chunks buffer qua WebSocket"),
        ("Transcription", "WhisperX Medium → text + timestamp"),
        ("Speaker Diarization", "PyAnnote 3.1 phân biệt người nói"),
        ("Lưu & Hiển thị", "Save DB, hiển thị real-time")
    ],
    tech_items=[
        ("WhisperX", "Speech-to-Text 92%, 50+ ngôn ngữ"),
        ("PyAnnote 3.1", "Speaker Diarization"),
        ("WebRTC + WebSocket", "Real-time streaming"),
        ("FFmpeg", "Audio processing")
    ],
    highlight="DiarizationPipeline wrapper đảm bảo format tương thích"
)

# Slide 4: Upload & Analysis
add_feature_slide(
    prs,
    title="Upload & Phân Tích",
    icon="📤",
    goal="Upload file và nhận phân tích AI toàn diện",
    stats=[
        ("100MB", "Max Size"),
        ("5", "Formats"),
        ("30s", "Avg Time")
    ],
    workflow_steps=[
        ("Upload File", "TXT, DOCX, MP3, WAV, MP4"),
        ("File Validation", "Check format, size, virus scan"),
        ("Content Extraction", "Audio→WhisperX, Text→Direct"),
        ("AI Analysis", "Gemini 2.5: tóm tắt, topics, actions"),
        ("Vectorization", "Embeddings → ChromaDB + PostgreSQL")
    ],
    tech_items=[
        ("Gemini 2.5 Flash", "LLM phân tích & trích xuất"),
        ("Prompt Engineering", "Few-shot, Chain-of-Thought"),
        ("python-docx", "Đọc/ghi Word"),
        ("Input Sanitization", "Validate, clean input")
    ],
    highlight="Structured Output với JSON schema đảm bảo format nhất quán"
)

# Slide 5: RAG Chat
add_feature_slide(
    prs,
    title="Chat với AI (RAG)",
    icon="💬",
    goal="Hỏi đáp thông minh với Retrieval-Augmented Generation",
    stats=[
        ("384", "Embed Dims"),
        ("5", "Top-K"),
        ("10", "History")
    ],
    workflow_steps=[
        ("User Query", "Người dùng đặt câu hỏi"),
        ("Query Embedding", "Câu hỏi → vector 384D"),
        ("Similarity Search", "Tìm top-k chunks (cosine)"),
        ("Context Assembly", "Retrieved + history"),
        ("LLM Generation", "Gemini tạo câu trả lời"),
        ("Memory Update", "Lưu Q&A vào memory")
    ],
    tech_items=[
        ("ChromaDB", "Vector Database"),
        ("LangChain", "RAG framework"),
        ("sentence-transformers", "Embeddings 384D"),
        ("Conversation Memory", "Follow-up questions")
    ],
    highlight="RAG + Memory = chatbot hiểu context và câu hỏi tiếp theo"
)

# Slide 6: Semantic Search
add_feature_slide(
    prs,
    title="Tìm Kiếm Ngữ Nghĩa",
    icon="🔍",
    goal="Tìm kiếm theo ý nghĩa, không chỉ từ khóa",
    stats=[
        ("384", "Dims"),
        ("95%", "Accuracy"),
        ("100ms", "Speed")
    ],
    workflow_steps=[
        ("Search Query", "User nhập query"),
        ("Query Embedding", "Query → vector 384D"),
        ("Vector Search", "ChromaDB + metadata filters"),
        ("Ranking", "Sắp xếp theo similarity (0-1)"),
        ("Display Results", "Hiển thị với highlight")
    ],
    tech_items=[
        ("ChromaDB Advanced", "Metadata filtering, search"),
        ("Cosine Similarity", "Đo tương đồng (0-1)"),
        ("Metadata Filters", "Ngày, ngôn ngữ, loại"),
        ("Analytics", "Thống kê database")
    ],
    highlight="Vector search + metadata filtering = tìm kiếm chính xác & nhanh"
)

# Save presentation
prs.save('Meeting_Chatbot_Presentation_Dark.pptx')
print("Da tao file: Meeting_Chatbot_Presentation_Dark.pptx")
